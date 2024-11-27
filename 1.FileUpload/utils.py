import os
import requests
from fastapi import UploadFile, HTTPException
from PyPDF2 import PdfReader  # For PDF handling
from pptx import Presentation  # For PowerPoint handling
from docx import Document  # For Word document handling

# Directory for saving uploaded files
UPLOAD_DIR = "./uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

# Allowed file formats
ALLOWED_FORMATS = {"pdf", "ppt", "doc"}

def validate_file(file: UploadFile):
    """Validate the uploaded file format."""
    file_format = file.filename.split(".")[-1].lower()
    if file_format not in ALLOWED_FORMATS:
        raise HTTPException(status_code=400, detail="File type not supported")

def save_file(file: UploadFile) -> str:
    """Save the uploaded file to the server."""
    file_path = os.path.join(UPLOAD_DIR, file.filename)
    with open(file_path, "wb") as f:
        data = file.file.read()
        f.write(data)
    return file_path

def extract_text(file_path: str) -> str:
    """Extract text based on file type."""
    if file_path.endswith(".pdf"):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith(".ppt"):
        return extract_text_from_ppt(file_path)
    elif file_path.endswith(".doc"):
        return extract_text_from_doc(file_path)
    else:
        raise HTTPException(status_code=400, detail="File type not supported")

def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from a PDF file."""
    reader = PdfReader(file_path)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_ppt(file_path: str) -> str:
    """Extract text from a PowerPoint file."""
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text + "\n"
    return text

def extract_text_from_doc(file_path: str) -> str:
    """Extract text from a Word document."""
    doc = Document(file_path)
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text

def send_processing(file_id: int, text: str):
    """Send the extracted text to the processing endpoint."""
    response = requests.post(
        "http://localhost:8001/process/",
        json={"file_id": file_id, "text": text}
    )
    if response.status_code != 200:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to send text for processing: {response.text}"
        )
